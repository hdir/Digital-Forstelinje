#!/usr/bin/env python3
"""
Generate PowerPoint presentation with barriers as bubble chart.
Each barrier is represented as a circle sized by number of actors.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import math

# Data from the table
barriers = [
    {"num": 1, "name": "Fragmenterte pasientforløp, svak samhandling og uklare ansvarslinjer", "actors": 101, "groups": 7},
    {"num": 2, "name": "Personellmangel, kapasitetskrise og kompetansegap", "actors": 97, "groups": 7},
    {"num": 3, "name": "Digital fragmentering og manglende interoperabilitet", "actors": 76, "groups": 6},
    {"num": 4, "name": "Finansieringsmodeller som gir feil insentiver (siloøkonomi)", "actors": 66, "groups": 6},
    {"num": 5, "name": "Forebygging, tidlig innsats og mestring underprioriteres", "actors": 58, "groups": 5},
    {"num": 6, "name": "Psykisk helse og rus er systematisk underprioritert", "actors": 32, "groups": 3},
    {"num": 7, "name": "Underutnyttelse av helsepersonells kompetanse og fastlåst oppgavedeling", "actors": 30, "groups": 3},
    {"num": 8, "name": "Geografisk og sosial ulikhet – «the missing middle» og distriktssvikt", "actors": 29, "groups": 4},
    {"num": 9, "name": "Pilot til drift: innovasjoner skaleres ikke", "actors": 20, "groups": 3},
    {"num": 10, "name": "Rehabilitering og habilitering er svekket og fragmentert", "actors": 19, "groups": 4},
    {"num": 11, "name": "Trege og lite innovasjonsvennlige anskaffelser og metodevurderinger", "actors": 17, "groups": 3},
    {"num": 12, "name": "Svekket primærhelsetjeneste og ubalansert oppgaveoverføring", "actors": 16, "groups": 2},
    {"num": 13, "name": "Svak forsknings- og innovasjonskapasitet i kommunene", "actors": 8, "groups": 3},
    {"num": 14, "name": "Tannhelse er ikke integrert i øvrig helsetjeneste", "actors": 6, "groups": 2},
    {"num": 15, "name": "Ideelle aktørers rammevilkår svekkes", "actors": 5, "groups": 2},
]

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Add blank slide
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

# Add title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
title_frame = title_box.text_frame
title_frame.text = "Gjentakende barrierer i Helsereformen"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 51, 102)

# Calculate bubble positions using a simple force-directed layout
# Normalize actor count to bubble size (in inches)
min_actors = min(b["actors"] for b in barriers)
max_actors = max(b["actors"] for b in barriers)

min_size = 0.78  # inches (0.6 * 1.3)
max_size = 2.6  # inches (2.0 * 1.3)

# Position bubbles in a grid-like pattern with some randomness
positions = []
cols = 4
rows = 4

import random
random.seed(42)

x_positions = [1.5, 3.5, 5.5, 7.5]
y_positions = [1.5, 2.8, 4.1, 5.4]

# Shuffle barriers for better visual distribution
shuffled_barriers = sorted(barriers, key=lambda x: -x["actors"])

for i, barrier in enumerate(shuffled_barriers):
    col = i % cols
    row = i // cols
    
    # Add some randomness to avoid perfect grid
    x = x_positions[col] + random.uniform(-0.3, 0.3)
    y = y_positions[row] + random.uniform(-0.2, 0.2)
    
    # Calculate bubble size based on actor count
    normalized = (barrier["actors"] - min_actors) / (max_actors - min_actors)
    size = min_size + normalized * (max_size - min_size)
    
    # Add circle (oval shape)
    left = Inches(x - size/2)
    top = Inches(y - size/2)
    width = Inches(size)
    height = Inches(size)
    
    # Color based on number of groups
    if barrier["groups"] == 7:
        color = RGBColor(220, 20, 60)  # Crimson - most critical
    elif barrier["groups"] >= 6:
        color = RGBColor(255, 69, 0)  # Red-orange
    elif barrier["groups"] >= 5:
        color = RGBColor(255, 140, 0)  # Dark orange
    elif barrier["groups"] >= 4:
        color = RGBColor(255, 165, 0)  # Orange
    elif barrier["groups"] >= 3:
        color = RGBColor(255, 200, 0)  # Gold
    else:
        color = RGBColor(255, 215, 0)  # Goldenrod
    
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    oval.fill.solid()
    oval.fill.fore_color.rgb = color
    oval.line.color.rgb = RGBColor(50, 50, 50)
    oval.line.width = Pt(1.5)
    
    # Add barrier name as text on the circle
    text_frame = oval.text_frame
    text_frame.word_wrap = True
    text_frame.clear()
    text_frame.margin_top = Inches(0.1)
    text_frame.margin_bottom = Inches(0.1)
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)
    text_frame.vertical_anchor = 1  # Middle
    
    # Add barrier name
    p1 = text_frame.paragraphs[0]
    p1.text = barrier['name']
    p1.font.size = Pt(10)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 255, 255)
    p1.alignment = 1  # Center
    
    positions.append((barrier, x, y))

# Add legend
legend_y = 6.5
legend_box = slide.shapes.add_textbox(Inches(0.5), Inches(legend_y), Inches(9), Inches(0.8))
legend_frame = legend_box.text_frame
legend_frame.word_wrap = True
legend_text = legend_frame.paragraphs[0]
legend_text.text = "Farger indikerer antall interessentgrupper som omtaler barrieren (fra rødt=7 grupper til gult=2 grupper). Størrelse indikerer antall aktører."
legend_text.font.size = Pt(9)
legend_text.font.italic = True
legend_text.font.color.rgb = RGBColor(100, 100, 100)

# Save presentation
output_path = r"c:\Git\Digital-Forstelinje\background\helsereform-august\Barriers_Bubble_Chart_Updated.pptx"
prs.save(output_path)
print(f"Presentation created: {output_path}")
