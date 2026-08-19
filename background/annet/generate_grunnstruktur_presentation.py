from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_PATH = Path(__file__).with_name("grunnstruktur-helsetjenesten.pptx")

NAVY = RGBColor(19, 54, 79)
INK = RGBColor(38, 50, 56)
WHITE = RGBColor(255, 255, 255)
MINT = RGBColor(220, 240, 232)
BLUE = RGBColor(218, 234, 247)
PEACH = RGBColor(252, 229, 213)
YELLOW = RGBColor(255, 244, 194)
GREY = RGBColor(236, 240, 242)
ARROW = RGBColor(74, 92, 101)


def add_box(slide, name, x, y, width, height, fill, font_size=13):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = NAVY
    shape.line.width = Pt(1.25)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.04)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.text = name
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.name = "Aptos"
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = True
    paragraph.font.color.rgb = INK
    return shape


def add_arrow(slide, start, end, label, label_x, label_y, color=ARROW, dashed=False):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(start[0]), Inches(start[1]), Inches(end[0]), Inches(end[1]),
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(1.5)
    connector.line.end_arrowhead = True
    if dashed:
        connector.line.dash_style = 1
    if label:
        text = slide.shapes.add_textbox(
            Inches(label_x), Inches(label_y), Inches(1.2), Inches(0.25)
        )
        text.text_frame.clear()
        paragraph = text.text_frame.paragraphs[0]
        paragraph.text = label
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(9)
        paragraph.font.color.rgb = color


def main():
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = WHITE

    title = slide.shapes.add_textbox(Inches(0.45), Inches(0.22), Inches(12.4), Inches(0.45))
    title.text_frame.clear()
    paragraph = title.text_frame.paragraphs[0]
    paragraph.text = "Grunnstrukturen i helsetenesta"
    paragraph.font.name = "Aptos Display"
    paragraph.font.size = Pt(26)
    paragraph.font.bold = True
    paragraph.font.color.rgb = NAVY

    subtitle = slide.shapes.add_textbox(Inches(0.48), Inches(0.72), Inches(12), Inches(0.3))
    subtitle.text_frame.clear()
    paragraph = subtitle.text_frame.paragraphs[0]
    paragraph.text = "Vanlege pasientforløp mellom kommunehelsetenesta og spesialisthelsetenesta"
    paragraph.font.name = "Aptos"
    paragraph.font.size = Pt(12)
    paragraph.font.color.rgb = RGBColor(85, 101, 108)

    zone_specs = [
        ("KOMMUNEHELSETENESTA", 0.35, 1.18, 4.1, 5.8, MINT),
        ("OVERGANG OG AKUTT", 4.62, 1.18, 3.2, 5.8, YELLOW),
        ("SPESIALISTHELSETENESTA", 7.98, 1.18, 5.0, 5.8, BLUE),
    ]
    for name, x, y, width, height, fill in zone_specs:
        zone = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
        zone.fill.solid()
        zone.fill.fore_color.rgb = fill
        zone.fill.transparency = 38
        zone.line.color.rgb = RGBColor(165, 181, 187)
        zone.line.width = Pt(0.75)
        heading = slide.shapes.add_textbox(Inches(x + 0.14), Inches(y + 0.12), Inches(width - 0.28), Inches(0.25))
        heading.text_frame.clear()
        paragraph = heading.text_frame.paragraphs[0]
        paragraph.text = name
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(10)
        paragraph.font.bold = True
        paragraph.font.color.rgb = NAVY

    boxes = {
        "patient": ("Pasient / innbyggjar", 0.75, 2.0, 3.3, 0.62, WHITE, 14),
        "gp": ("Fastlege", 0.75, 3.0, 1.55, 0.68, WHITE, 14),
        "municipal": ("Kommunale helse- og\nomsorgstenester", 2.5, 2.9, 1.55, 0.88, WHITE, 11),
        "rehab": ("Rehabilitering\nog opptrening", 2.5, 4.65, 1.55, 0.78, WHITE, 11),
        "emergency": ("Legevakt\n113 / 116 117", 4.98, 2.0, 1.45, 0.78, PEACH, 11),
        "ambulance": ("Ambulansetenesta", 4.98, 4.2, 1.45, 0.68, PEACH, 11),
        "dms": ("Distriktsmedisinsk\nsenter", 6.45, 3.05, 1.15, 0.88, PEACH, 10),
        "specialist": ("Spesialist /\npoliklinikk", 8.35, 2.0, 1.5, 0.78, WHITE, 11),
        "dps": ("Distriktspsykiatrisk\nsenter", 10.1, 2.0, 1.8, 0.78, WHITE, 11),
        "hospital": ("Sjukehus", 9.05, 4.15, 1.7, 0.7, WHITE, 14),
        "regional": ("Regionale helseføretak\n(eigar og ansvar)", 11.0, 5.45, 1.7, 0.8, GREY, 10),
    }

    add_arrow(slide, (2.4, 2.31), (2.5, 3.34), "kontakt", 1.55, 2.55)
    add_arrow(slide, (2.3, 3.34), (2.5, 3.34), "samarbeid", 2.05, 3.12)
    add_arrow(slide, (2.3, 3.34), (2.5, 5.04), "oppfølging", 1.6, 4.18)
    add_arrow(slide, (4.05, 3.2), (4.98, 2.39), "akutt hjelp", 4.0, 2.35)
    add_arrow(slide, (4.05, 3.2), (6.45, 3.49), "samlokalisert", 4.55, 3.05)
    add_arrow(slide, (4.05, 3.2), (8.35, 2.39), "tilvising", 6.0, 2.23)
    add_arrow(slide, (8.35, 2.39), (10.1, 2.39), "tilvising", 8.85, 2.12)
    add_arrow(slide, (8.35, 2.39), (9.9, 4.15), "innlegging", 8.25, 3.25)
    add_arrow(slide, (10.1, 2.39), (9.9, 4.15), "meir spesialisert", 10.05, 3.28)
    add_arrow(slide, (5.7, 2.78), (9.05, 4.15), "akutt transport", 6.6, 3.25)
    add_arrow(slide, (9.9, 4.85), (4.05, 5.04), "tilbake til kommunen", 6.2, 4.86)
    add_arrow(slide, (4.05, 5.04), (4.05, 3.7), "vidare oppfølging", 3.55, 4.0, dashed=True)
    add_arrow(slide, (11.85, 5.45), (10.75, 4.85), "styrer/eiger", 11.05, 5.0, dashed=True)

    for name, (label, x, y, width, height, fill, font_size) in boxes.items():
        add_box(slide, label, x, y, width, height, fill, font_size)

    footer = slide.shapes.add_textbox(Inches(0.48), Inches(7.12), Inches(12.3), Inches(0.2))
    footer.text_frame.clear()
    paragraph = footer.text_frame.paragraphs[0]
    paragraph.text = "Kjelde: Helse- og omsorgsdepartementet, «Grunnstrukturen i helsetenesta», sist oppdatert 07.05.2025"
    paragraph.font.name = "Aptos"
    paragraph.font.size = Pt(8)
    paragraph.font.color.rgb = RGBColor(105, 117, 122)

    presentation.save(OUTPUT_PATH)
    print(OUTPUT_PATH)


if __name__ == "__main__":
    main()